import os, logging, sys, io
import win32com.client as win32
import win32clipboard
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.util import Pt, Inches
import openpyxl

ppLayoutBlank = 12  # Define the default layout index for a blank slide

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def copy_chart_to_existing_slide(excel_path, sheet_name, presentation, slide_index):
    xlApp = win32.Dispatch('Excel.Application')

    try:
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = True

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()
                copiedSheet.ChartObjects(1).CopyPicture()

                # Open the clipboard and get clipboard data
                win32clipboard.OpenClipboard(0)
                clipboard_data = win32clipboard.GetClipboardData(win32clipboard.CF_ENHMETAFILE)
                win32clipboard.CloseClipboard()

                # Create a BytesIO object to handle the clipboard data
                image_stream = io.BytesIO(clipboard_data)

                # Get the existing slide by index
                slide = presentation.slides[slide_index]
                shapes = slide.shapes

                # Define the area to cover on the slide
                if slide_index == 9:  # For MRR CHART
                    left = Inches(5.2)
                    top = Inches(1.3)
                    width = Inches(8.2)
                    height = Inches(6.5)

                else: # For CSR and CHURCH ANALYSIS CHARTS
                    left = Inches(1.5)
                    top = Inches(0.4)
                    width = Inches(11.5)
                    height = Inches(6.7)          

                # Add the picture from the BytesIO object with the specified area
                shapes.add_picture(image_stream, left, top, width=width, height=height)

                return

        logging.warning(f"No chart detected in sheet: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def copy_avg_attendance_to_slide(excel_path, sheet_name, presentation, slide_index, copy_count=None):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                # Get data from the sheet, excluding the first column (column A)
                data = [list(row)[1:] for row in copiedSheet.UsedRange.Value]

                # Create a new slide in the PowerPoint presentation
                slide = presentation.slides[slide_index]

                # Define the area to cover on the slide based on slide_index
                if slide_index == 4:  # Slide 4 - First 16 rows
                    left = Inches(1.5)
                    top = Inches(0.5)
                    width = Inches(13)
                    height = Inches(6.8)

                    # Determine the number of rows to copy based on copy_count
                    if copy_count:
                        data = data[:copy_count]

                elif slide_index == 5:  # Slide 5 - First row + last 15 rows of column A
                    left = Inches(1.5)
                    top = Inches(0.5)
                    width = Inches(13)
                    height = Inches(6.9)

                    # Extract the first row and bottom 15 rows of column A excluding the last row
                    data = [data[0]] + data[-16:-1]

                # Add data to the slide, for example, create a table
                table = slide.shapes.add_table(rows=len(data), cols=len(data[0]), left=left, top=top,
                                              width=width, height=height).table

                # Populate the table with data
                for i, row in enumerate(data):
                    for j, value in enumerate(row):
                        cell = table.cell(i, j)
                        cell.text = str(value)

                        # Adjust font size if specified
                        if copy_count and i == copy_count - 1:
                            cell.text_frame.paragraphs[0].font.size = Pt(11)  
                        else:
                            # Default font size for other cells
                            cell.text_frame.paragraphs[0].font.size = Pt(11)

                return

        logging.warning(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def copy_table_to_slide(excel_path, sheet_name, presentation, slide_index, font_size=None, copy_count=None):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                # Get data from the sheet, assuming you want to copy cell values
                data = copiedSheet.UsedRange.Value

                # Determine the number of rows to copy based on copy_count
                if copy_count:
                    data = data[:copy_count]

                # Create a new slide in the PowerPoint presentation
                slide = presentation.slides[slide_index]

                # Define the area to cover on the slide based on slide_index
                if slide_index == 12:  # CSR Table 2
                    left = Inches(1.5)
                    top = Inches(1.3)
                    width = Inches(8.2)
                    height = Inches(6.5)
                elif slide_index == 8:  # MRR
                    left = Inches(1.5)
                    top = Inches(0.4)
                    width = Inches(10)
                    height = Inches(6)
                elif slide_index == 11:  # CSR Table 1
                    left = Inches(1.5)
                    top = Inches(0.4)
                    width = Inches(10)
                    height = Inches(6)
                elif slide_index == 22:  # Church Analysis
                    left = Inches(1.5)
                    top = Inches(0.4)
                    width = Inches(11.6)
                    height = Inches(6.5)
                else:  # HF, Converts
                    data = [list(row)[1:] for row in copiedSheet.UsedRange.Value]
                    left = Inches(1.5)
                    top = Inches(0.4)
                    width = Inches(11.5)
                    height = Inches(6.7)

                # Add data to the slide, for example, create a table
                table = slide.shapes.add_table(rows=len(data), cols=len(data[0]), left=left, top=top,
                                              width=width, height=height).table

                # Populate the table with data
                for i, row in enumerate(data):
                    for j, value in enumerate(row):
                        cell = table.cell(i, j)
                        cell.text = str(value)

                        # Adjust font size if specified
                        if font_size:
                            cell.text_frame.paragraphs[0].font.size = font_size

                return

        logging.warning(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def process_excel_files_in_folder(folder_path, pptx_template_path):
    xlApp = win32.Dispatch('Excel.Application')
    pptApp = win32.Dispatch('PowerPoint.Application')

    try:
        # Load the PowerPoint template once outside the loop
        template_presentation = Presentation(pptx_template_path)

        for file_name in os.listdir(folder_path):
            if file_name.endswith('.xlsx'):
                # Reuse the same template for each Excel file
                presentation = Presentation(pptx_template_path)

                excel_path = os.path.join(folder_path, file_name)
                excel_file_name = os.path.splitext(file_name)[0]

                # Instructions to copy chart slide
                copy_chart_to_existing_slide(excel_path, 'MRR', presentation, 9)
                copy_chart_to_existing_slide(excel_path, 'CSR DISTRIBUTION', presentation, 13)
                copy_chart_to_existing_slide(excel_path, 'CHURCH ANALYSIS1', presentation, 23)

                # Instructions to copy table to slide
                copy_avg_attendance_to_slide(excel_path, 'AVG. ATTENDANCE', presentation, 4, copy_count=16)
                copy_avg_attendance_to_slide(excel_path, 'AVG. ATTENDANCE', presentation, 5, copy_count=16)
                copy_table_to_slide(excel_path, 'MRR', presentation, 8, font_size=Pt(15))
                copy_table_to_slide(excel_path, 'CSR', presentation, 11, font_size=Pt(15))
                copy_table_to_slide(excel_path, 'CSR DISTRIBUTION', presentation, 12, font_size=Pt(15))
                copy_table_to_slide(excel_path, 'CONVERTS', presentation, 16, copy_count=16, font_size=Pt(7))
                copy_table_to_slide(excel_path, 'HOUSE FELLOWSHIP', presentation, 19, copy_count=16, font_size=Pt(11))
                copy_table_to_slide(excel_path, 'CHURCH ANALYSIS1', presentation, 22, copy_count=16, font_size=Pt(15))

                output_pptx_path = os.path.join(folder_path, f'{excel_file_name}_updated.pptx')
                presentation.save(output_pptx_path)

                logging.info(f'Finished populating {file_name}')

    except Exception as e:
        logging.error(f"Error during processing: {e}")

    finally:
        # Quit Excel and PowerPoint after processing all files
        xlApp.Quit()
        pptApp.Quit()

def main():
    if len(sys.argv) == 3:
        folder_path = sys.argv[1]
        pptx_path = sys.argv[2]
    else:
        # Use default paths if command-line arguments are not provided
        folder_path = r"C:\Users\PMD - FEMI\Desktop\cleanedFiles\cleaned_files"
        pptx_path = r'C:\Users\PMD - FEMI\Desktop\region 13\sampleAutomate.pptx'

    process_excel_files_in_folder(folder_path, pptx_path)

if __name__ == "__main__":
    main()


