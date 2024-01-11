import os, logging, sys, io, time
import win32com.client as win32
import win32com.client
import win32clipboard
from pptx.dml.color import RGBColor
from pptx import Presentation
from win32com.client import Dispatch
from pptx.util import Pt
from pptx.util import Inches
from pptx.util import Pt, Inches
import openpyxl

ppLayoutBlank = 12  # Define the default layout index for a blank slide

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
"""
def copy_chart_to_existing_slide(excel_path, sheet_name, presentation, slide_index):
    xlApp = win32.Dispatch('Excel.Application')

    try:
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = True

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()
                copiedSheet.ChartObjects(1).CopyPicture()

                # Open the clipboard and get clipboard data
                win32clipboard.OpenClipboard(0)
                clipboard_data = win32clipboard.GetClipboardData(win32clipboard.CF_ENHMETAFILE)
                win32clipboard.CloseClipboard()

                # Create a BytesIO object to handle the clipboard data
                image_stream = io.BytesIO(clipboard_data)

                # Get the existing slide by index
                slide = presentation.slides[slide_index]
                shapes = slide.shapes

                # Define the area to cover on the slide
                if slide_index == 9:  # For MRR CHART
                    left = Inches(5.2)
                    top = Inches(1.3)
                    width = Inches(8.2)
                    height = Inches(6.5)

                else: # For CSR and CHURCH ANALYSIS CHARTS
                    left = Inches(1.5)
                    top = Inches(0.4)
                    width = Inches(11.5)
                    height = Inches(6.7)          

                # Add the picture from the BytesIO object with the specified area
                shapes.add_picture(image_stream, left, top, width=width, height=height)

                return

        logging.warning(f"No chart detected in sheet: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)
"""
def copy_avg_attendance_to_slide(excel_path, sheet_name, presentation, slide_index, copy_count=None):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                # Get data from the sheet, excluding the first column (column A)
                data = [list(row)[1:] for row in copiedSheet.UsedRange.Value]

                # Round each value to the nearest whole number
                data = [[round(cell) if isinstance(cell, (float, int)) else cell for cell in row] for row in data]

                # Create a new slide in the PowerPoint presentation
                slide = presentation.slides[slide_index]

                # Define the area to cover on the slide based on slide_index
                if slide_index == 4:  # Slide 4 - First 16 rows
                    left = Inches(1.5)
                    top = Inches(0.2)
                    width = Inches(11.6)
                    height = Inches(5.2)

                    # Determine the number of rows to copy based on copy_count
                    if copy_count:
                        data = data[:copy_count]

                elif slide_index == 5:  # Slide 5 - First row + last 15 rows of column A
                    left = Inches(1.5)
                    top = Inches(0.2)
                    width = Inches(11.7)
                    height = Inches(5.8)

                    # Extract the first row and bottom 15 rows of column A excluding the last row
                    data = [data[0]] + data[-16:-1]

                # Add data to the slide, for example, create a table
                table = slide.shapes.add_table(rows=len(data), cols=len(data[0]), left=left, top=top,
                                              width=width, height=height).table

                # Set the background color of the first row to a purple accent
                for j in range(len(data[0])):
                    cell = table.cell(0, j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(128, 0, 128)  # Purple Accent
                
                # Populate the table with data
                for i, row in enumerate(data):
                    for j, value in enumerate(row):
                        cell = table.cell(i, j)
                        cell.text = str(value)

                        # Adjust font size if specified
                        if copy_count and i == copy_count - 1:
                            cell.text_frame.paragraphs[0].font.size = Pt(13)  
                        else:
                            # Default font size for other cells
                            cell.text_frame.paragraphs[0].font.size = Pt(12)

                return

        logging.warning(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def copy_table_to_slide(excel_path, sheet_name, presentation, slide_index, font_size=None, copy_count=None):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                # Define default data extraction method
                data_extraction_method = copiedSheet.UsedRange.Value

                # Override data extraction for specific sheets
                if sheet_name in ["HOUSE FELLOWSHIP", "CONVERTS"]:
                    data_extraction_method = [list(row)[1:] for row in copiedSheet.UsedRange.Value]

                # Get data from the sheet
                data = data_extraction_method

                # Convert numeric values to integers
                #data = [[int(cell) if isinstance(cell, (float, int)) else cell for cell in row] for row in data]

                # Determine the number of rows to copy based on copy_count
                if copy_count:
                    data = data[:copy_count]

                # Round each numeric value to the nearest whole number, excluding "CHURCH ANALYSIS1" sheet
                if sheet_name != "CHURCH ANALYSIS1":
                    data = [[round(cell) if isinstance(cell, (float, int)) else cell for cell in row] for row in data]

                # Create a new slide in the PowerPoint presentation
                slide = presentation.slides[slide_index]

                # Define the area to cover on the slide based on slide_index
                if slide_index == 12:  # CSR Distribution
                    left = Inches(1.5)
                    top = Inches(0.5)
                    width = Inches(11.5)
                    height = Inches(5.5)
                elif slide_index == 8:  # MRR
                    left = Inches(1.3)
                    top = Inches(0.7)
                    width = Inches(11.8)
                    height = Inches(6.2)
                elif slide_index == 11:  # CSR Table 1
                    left = Inches(1.6)
                    top = Inches(0.3)
                    width = Inches(11.5)
                    height = Inches(7)
                elif slide_index == 22:  # Church Analysis
                    left = Inches(1.5)
                    top = Inches(0.4)
                    width = Inches(11.6)
                    height = Inches(6.5)
                elif slide_index == 19:  # House Fellowship
                    left = Inches(1.5)
                    top = Inches(0.5)
                    width = Inches(11.5)
                    height = Inches(6.7)
                elif slide_index == 16:  # Converts
                    left = Inches(1.4)
                    top = Inches(0.5)
                    width = Inches(11.6)
                    height = Inches(6.7)
                else:  # Default
                    left = Inches(1.5)
                    top = Inches(0.4)
                    width = Inches(11.5)
                    height = Inches(6.7)

                # Add data to the slide, for example, create a table
                table = slide.shapes.add_table(rows=len(data), cols=len(data[0]), left=left, top=top,
                                              width=width, height=height).table

                # Set the background color of the first row to a purple accent
                for j in range(len(data[0])):
                    cell = table.cell(0, j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(128, 0, 128)  # Purple Accent
                
                
                # Populate the table with data
                for i, row in enumerate(data):
                    for j, value in enumerate(row):
                        cell = table.cell(i, j)
                        
                        try:
                            # Attempt to convert the value to a string and set it in the cell
                            cell.text = str(value)
                        except Exception as e:
                            # If there's an exception (e.g., trying to convert 'None' to string), print an error message
                            print(f"Error setting cell ({i}, {j}) value: {e}")            
                
               # Check if the sheet name is "CHURCH ANALYSIS"
                if sheet_name == "CHURCH ANALYSIS1":
                    # Multiply the values in the second column by 100
                    for i in range(1, len(data)):
                        cell = table.cell(i, 2)  # Assuming the second column (index 2)
                        value = float(cell.text)  # Convert the text to a float for multiplication
                        cell.text = f"{value * 100:.2f}%"               

                return

        logging.warning(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def copy_cells_to_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                # Find the row index where "TOTAL" is present in column B
                total_row = None
                for row in range(1, copiedSheet.UsedRange.Rows.Count + 1):
                    if copiedSheet.Cells(row, 2).Value == "TOTAL":
                        total_row = row
                        break

                if total_row is None:
                    raise ValueError("Row with 'TOTAL' not found in column B.")

                # Define the columns to extract based on sheet name
                if sheet_name == 'AVG. ATTENDANCE':
                    column_indices = [3, 7, 8, 9]
                else:
                    column_indices = [3, 7, 8]

                # Get values from specified cells in the "TOTAL" row based on sheet name
                numeric_values = [copiedSheet.Cells(total_row, col).Value for col in column_indices]
                cell_values = []

                for value in numeric_values:
                    if isinstance(value, (int, float)):
                        cell_values.append(round(value))
                    else:
                        # Handle non-numeric values, such as strings
                        cell_values.append(value)

                # Create a new slide in the PowerPoint presentation
                slide = presentation.slides[slide_index]

                # Add data to the slide with individual positions for each textbox:left,top,width,height
                textbox1 = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(1), Inches(0.5))
                textbox1.text_frame.text = f"{cell_values[0]}"

                textbox2 = slide.shapes.add_textbox(Inches(6.5), Inches(3.7), Inches(1), Inches(0.5))
                textbox2.text_frame.text = f"{cell_values[1]}"

                textbox3 = slide.shapes.add_textbox(Inches(11), Inches(5.6), Inches(1), Inches(0.5))
                textbox3.text_frame.text = f"{cell_values[2]}"
                
                if sheet_name == 'AVG. ATTENDANCE':
                    textbox4 = slide.shapes.add_textbox(Inches(11.3), Inches(4.5), Inches(1), Inches(0.5))
                    textbox4.text_frame.text = f"{cell_values[3]}"

                # Add more text boxes as needed

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def copy_csrcells_to_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                if sheet_name == 'CSR':
                    # Find the column index where "Nov 2023" is present in row 1 of column D
                    total_column = None
                    for col in range(1, copiedSheet.UsedRange.Columns.Count + 1):
                        if copiedSheet.Cells(1, col).Value == 'TOTAL':
                            total_column = col
                            break

                    if total_column is not None:
                        #cell_value = [copiedSheet.Cells(3, col).Value for col in range(nov_2023_column, nov_2023_column + 1)]
                        column = total_column
                        csr_pro = round(copiedSheet.Cells(2, total_column).Value)
                        csr_fin = round(copiedSheet.Cells(3, total_column).Value)
                        csr_ben = round(copiedSheet.Cells(5, total_column).Value)
                        #csr_value = round(copiedSheet.Cells(2, total_column).Value)

                        # Create a new slide in the PowerPoint presentation
                        slide = presentation.slides[slide_index]

                        # Add data to the slide with individual positions for each textbox:left,top,width,height
                        textbox1 = slide.shapes.add_textbox(Inches(9.2), Inches(5.4), Inches(1), Inches(0.5))
                        textbox1.text_frame.text = f"{csr_pro}"
                       
                        textbox2 = slide.shapes.add_textbox(Inches(9), Inches(6.4), Inches(1), Inches(0.5))
                        textbox2.text_frame.text = f"{csr_fin}"

                        textbox3 = slide.shapes.add_textbox(Inches(9), Inches(2.1), Inches(1), Inches(0.5))
                        textbox3.text_frame.text = f"{csr_ben}"

                        #textbox4 = slide.shapes.add_textbox(Inches(3.2), Inches(3.7), Inches(1), Inches(0.5))
                        #textbox4.text_frame.text = f"{csr_value}"

                        # Add more text boxes as needed

                    else:
                        raise ValueError("Column 'Total' not found in sheet 'CSR'")

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def copy_mrrcells_to_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                if sheet_name == 'MRR':
                    # Find the column index where "Nov 2023" is present in row 1 of column D
                    nov_2023_column = None
                    for col in range(1, copiedSheet.UsedRange.Columns.Count + 1):
                        if copiedSheet.Cells(1, col).Value == 'Nov 2023':
                            nov_2023_column = col
                            break

                    if nov_2023_column is not None:
                        # Calculate MRR value and directly use it on the slide
                        mrr_value = round((copiedSheet.Cells(5, nov_2023_column).Value / 
                                           copiedSheet.Cells(3, nov_2023_column).Value) * 100, )
                        
                        #cell_value = [copiedSheet.Cells(3, col).Value for col in range(nov_2023_column, nov_2023_column + 1)]
                        column = nov_2023_column
                        first_timers_value = round(copiedSheet.Cells(3, nov_2023_column).Value)
                        avg_attendance_value = round(copiedSheet.Cells(4, nov_2023_column).Value)

                        # Create a new slide in the PowerPoint presentation
                        slide = presentation.slides[slide_index]

                        # Add data to the slide with individual positions for each textbox:left,top,width,height
                        # For MRR %
                        textbox1 = slide.shapes.add_textbox(Inches(4), Inches(5.7), Inches(1), Inches(0.5))
                        textbox1.text_frame.text = f"{mrr_value}%"
                        # For total cummulative value
                        textbox2 = slide.shapes.add_textbox(Inches(3.2), Inches(3.7), Inches(1), Inches(0.5))
                        textbox2.text_frame.text = f"{first_timers_value}"
                        #For avg. attendance
                        textbox2 = slide.shapes.add_textbox(Inches(3.4), Inches(1.5), Inches(1), Inches(0.5))
                        textbox2.text_frame.text = f"{avg_attendance_value}"

                        # Add more text boxes as needed

                    else:
                        raise ValueError("Column 'Nov 2023' not found in sheet 'MRR'")

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def copy_cacells_to_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                # Find the column index where "ATTENDANCE RANGE" is present in row 1
                attendance_range_column = None
                for col in range(1, copiedSheet.UsedRange.Columns.Count + 1):
                    if copiedSheet.Cells(1, col).Value == 'ATTENDANCE RANGE':
                        attendance_range_column = col
                        break

                if attendance_range_column is not None:
                    # Find the row index where "TOTAL" is present in the "ATTENDANCE RANGE" column
                    total_row = None
                    for row in range(1, copiedSheet.UsedRange.Rows.Count + 1):
                        cell_value = copiedSheet.Cells(row, attendance_range_column).Value
                        if cell_value == 'TOTAL':
                            total_row = row
                            break

                    # Calculate totalnoparishes_value based on the value in the next column of the "TOTAL" row
                    if total_row is not None:
                        totalnoparishes_value = int(copiedSheet.Cells(total_row, attendance_range_column + 1).Value)
                    else:
                        totalnoparishes_value = 0

                    # Find the row indices for the specified attendance ranges
                    range_50_99_row = None
                    range_10_49_row = None
                    range_1_9_row = None

                    for row in range(1, copiedSheet.UsedRange.Rows.Count + 1):
                        cell_value = copiedSheet.Cells(row, attendance_range_column).Value
                        if cell_value == '50 - 99 people':
                            range_50_99_row = row
                        elif cell_value == '10 - 49 people':
                            range_10_49_row = row
                        elif cell_value == '1 - 9 people':
                            range_1_9_row = row

                    # Calculate lt100percent_value based on the corresponding cells in the '% OF PARISHES' column
                    if range_50_99_row is not None or range_10_49_row is not None or range_1_9_row is not None:
                        total_50_99 = copiedSheet.Cells(range_50_99_row, copiedSheet.UsedRange.Columns.Count).Value if range_50_99_row is not None else 0
                        total_10_49 = copiedSheet.Cells(range_10_49_row, copiedSheet.UsedRange.Columns.Count).Value if range_10_49_row is not None else 0
                        total_1_9 = copiedSheet.Cells(range_1_9_row, copiedSheet.UsedRange.Columns.Count).Value if range_1_9_row is not None else 0

                        # Copy corresponding values from 'TOTAL NO. OF PARISHES' column
                        total_50_99_parishes = copiedSheet.Cells(range_50_99_row, copiedSheet.UsedRange.Columns.Count - 1).Value if range_50_99_row is not None else 0
                        total_10_49_parishes = copiedSheet.Cells(range_10_49_row, copiedSheet.UsedRange.Columns.Count - 1).Value if range_10_49_row is not None else 0
                        total_1_9_parishes = copiedSheet.Cells(range_1_9_row, copiedSheet.UsedRange.Columns.Count - 1).Value if range_1_9_row is not None else 0

                        # Calculate totalLT100_value based on the copied values
                        totalLT100_value = int(total_50_99_parishes + total_10_49_parishes + total_1_9_parishes)

                        lt100percent_value = round((total_50_99 + total_10_49 + total_1_9) * 100)
                    else:
                        lt100percent_value = 0
                        totalLT100_value = 0
                else:
                    raise ValueError("Column 'ATTENDANCE RANGE' not found in sheet 'CHURCH ANALYSIS1'")

                # Create a new slide in the PowerPoint presentation
                slide = presentation.slides[slide_index]

                # Add data to the slide with individual positions for each textbox:left,top,width,height
                # For percentage total number of parishes less than 100 within a province
                textbox1 = slide.shapes.add_textbox(Inches(5.9), Inches(5.5), Inches(1), Inches(0.5))
                textbox1.text_frame.text = f"{lt100percent_value}%"

                # For total number of parishes less than 100 within a province
                textbox2 = slide.shapes.add_textbox(Inches(5.3), Inches(5), Inches(1), Inches(0.5))
                textbox2.text_frame.text = f"{totalLT100_value}"

                # For total number of parishes (TOTAL)
                textbox3 = slide.shapes.add_textbox(Inches(1), Inches(5.1), Inches(1), Inches(0.5))
                textbox3.text_frame.text = f"{totalnoparishes_value}"

                # Add more text boxes as needed

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def recommendation_CA_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                # Find the column index where "ATTENDANCE RANGE" is present in row 1
                attendance_range_column = None
                for col in range(1, copiedSheet.UsedRange.Columns.Count + 1):
                    if copiedSheet.Cells(1, col).Value == 'ATTENDANCE RANGE':
                        attendance_range_column = col
                        break

                # Check if the column was found
                if attendance_range_column is not None:
                    # Find the row indices for the specified attendance ranges
                    range_50_99_row = None
                    range_10_49_row = None
                    range_1_9_row = None

                    for row in range(1, copiedSheet.UsedRange.Rows.Count + 1):
                        cell_value = copiedSheet.Cells(row, attendance_range_column).Value
                        if cell_value == '50 - 99 people':
                            range_50_99_row = row
                        elif cell_value == '10 - 49 people':
                            range_10_49_row = row
                        elif cell_value == '1 - 9 people':
                            range_1_9_row = row

                    # Calculate lt100percent_value based on the corresponding cells in the '% OF PARISHES' column
                    if range_50_99_row is not None or range_10_49_row is not None or range_1_9_row is not None:
                        total_50_99 = copiedSheet.Cells(range_50_99_row, copiedSheet.UsedRange.Columns.Count).Value if range_50_99_row is not None else 0
                        total_10_49 = copiedSheet.Cells(range_10_49_row, copiedSheet.UsedRange.Columns.Count).Value if range_10_49_row is not None else 0
                        total_1_9 = copiedSheet.Cells(range_1_9_row, copiedSheet.UsedRange.Columns.Count).Value if range_1_9_row is not None else 0

                        lt100percent_value = round((total_50_99 + total_10_49 + total_1_9) * 100)
                    else:
                        lt100percent_value = 0
                else:
                    raise ValueError("Column 'ATTENDANCE RANGE' not found in sheet 'CHURCH ANALYSIS1'")

                # Create a new slide in the PowerPoint presentation
                slide = presentation.slides[slide_index]

                # Add data to the slide with individual positions for each textbox:left,top,width,height
                # For percentage total number of parishes less than 100 within a province
                textbox1 = slide.shapes.add_textbox(Inches(2.4), Inches(4.7), Inches(1), Inches(0.5))
                textbox1.text_frame.text = f"{lt100percent_value}%"
                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def recommendation_MRR_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                if sheet_name == 'MRR':
                    # Find the column index where "Nov 2023" is present in row 1 of column D
                    nov_2023_column = None
                    for col in range(1, copiedSheet.UsedRange.Columns.Count + 1):
                        if copiedSheet.Cells(1, col).Value == 'Nov 2023':
                            nov_2023_column = col
                            break

                    if nov_2023_column is not None:
                        # Calculate MRR value and directly use it on the slide
                        mrr_value = round((copiedSheet.Cells(5, nov_2023_column).Value / 
                                           copiedSheet.Cells(3, nov_2023_column).Value) * 100, )

                        # Create a new slide in the PowerPoint presentation
                        slide = presentation.slides[slide_index]

                        # Add data to the slide with individual positions for each textbox:left,top,width,height
                        # For MRR %
                        textbox1 = slide.shapes.add_textbox(Inches(2.5), Inches(4.6), Inches(1), Inches(0.5))
                        textbox1.text_frame.text = f"{mrr_value}%"
                        
                    else:
                        raise ValueError("Column 'Nov 2023' not found in sheet 'MRR'")

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def recommendation_CSR_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()

                if sheet_name == 'CSR':
                    # Find the column index where "Nov 2023" is present in row 1 of column D
                    total_column = None
                    for col in range(1, copiedSheet.UsedRange.Columns.Count + 1):
                        if copiedSheet.Cells(1, col).Value == 'TOTAL':
                            total_column = col
                            break

                    if total_column is not None:
                        #cell_value = [copiedSheet.Cells(3, col).Value for col in range(nov_2023_column, nov_2023_column + 1)]
                        column = total_column                      
                        csr_fin = round(copiedSheet.Cells(3, total_column).Value)
                    
                        # Create a new slide in the PowerPoint presentation
                        slide = presentation.slides[slide_index]
                        # Add data to the slide with individual positions for each textbox:left,top,width,height                       
                        textbox2 = slide.shapes.add_textbox(Inches(2.1), Inches(2.3), Inches(1), Inches(0.5))
                        textbox2.text_frame.text = f"N{csr_fin}"

                    else:
                        raise ValueError("Column 'Total' not found in sheet 'CSR'")

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def recommendation_HF_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()
                
                # Define the columns to extract based on sheet name
                if sheet_name == 'HOUSE FELLOWSHIP':
                    # Find the row index where "TOTAL" is present in column B
                    total_row = None
                    for row in range(1, copiedSheet.UsedRange.Rows.Count + 1):
                        if copiedSheet.Cells(row, 2).Value == "TOTAL":
                            total_row = row
                            break

                    if total_row is None:
                        raise ValueError("Row with 'TOTAL' not found in column B.")
                    
                    # Define the columns to extract based on sheet name
                    column_indices = [3, 7]

                    # Get values from specified cells in the "TOTAL" row based on sheet name
                    numeric_values = [copiedSheet.Cells(total_row, col).Value for col in column_indices]
                    cell_values = []

                    for value in numeric_values:
                        if isinstance(value, (int, float)):
                            cell_values.append(round(value))
                        else:
                            # Handle non-numeric values, such as strings
                            cell_values.append(value)

                    # Calculate the desired value
                    result_value = ((cell_values[1] - cell_values[0]) / cell_values[0]) * 100

                    # Create a new slide in the PowerPoint presentation
                    slide = presentation.slides[slide_index]

                    # Add data to the slide with individual positions for each textbox:left,top,width,height
                    textbox1 = slide.shapes.add_textbox(Inches(2.3), Inches(1.9), Inches(1), Inches(0.5))
                    textbox1.text_frame.text = f"{round(result_value, 2)}%"  # Round to 2 decimal places 

                else:
                    raise ValueError("Column 'Total' not found in sheet 'CSR'")

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def recommendation_CONVERTS_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()
                
                # Define the columns to extract based on sheet name
                if sheet_name == 'CONVERTS':
                    # Find the row index where "TOTAL" is present in column B
                    total_row = None
                    for row in range(1, copiedSheet.UsedRange.Rows.Count + 1):
                        if copiedSheet.Cells(row, 2).Value == "TOTAL":
                            total_row = row
                            break

                    if total_row is None:
                        raise ValueError("Row with 'TOTAL' not found in column B.")
                    
                    # Define the columns to extract based on sheet name
                    column_indices = [3, 7]

                    # Get values from specified cells in the "TOTAL" row based on sheet name
                    numeric_values = [copiedSheet.Cells(total_row, col).Value for col in column_indices]
                    cell_values = []

                    for value in numeric_values:
                        if isinstance(value, (int, float)):
                            cell_values.append(round(value))
                        else:
                            # Handle non-numeric values, such as strings
                            cell_values.append(value)

                    # Calculate the desired value
                    result_value = ((cell_values[1] - cell_values[0]) / cell_values[0]) * 100

                    # Create a new slide in the PowerPoint presentation
                    slide = presentation.slides[slide_index]

                    # Add data to the slide with individual positions for each textbox:left,top,width,height
                    textbox1 = slide.shapes.add_textbox(Inches(2.3), Inches(4.7), Inches(1), Inches(0.5))
                    textbox1.text_frame.text = f"{round(result_value, 2)}%"  # Round to 2 decimal places 

                else:
                    raise ValueError("Column 'Total' not found in sheet 'CONVERTS'")

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def recommendation_ATTENDANCE_slide(excel_path, sheet_name, presentation, slide_index):
    try:
        xlApp = win32.Dispatch('Excel.Application')
        excelWorkbook = xlApp.Workbooks.Open(excel_path)

        # Find the specified sheet in the Excel workbook
        for sheet in excelWorkbook.Sheets:
            if sheet.Name == sheet_name:
                sheet.Copy(Before=excelWorkbook.Sheets[1])
                xlApp.Visible = False  # You may choose to hide Excel

                copiedSheet = excelWorkbook.Sheets[1]
                copiedSheet.Activate()
                
                # Define the columns to extract based on sheet name
                if sheet_name == 'AVG. ATTENDANCE':
                    # Find the row index where "TOTAL" is present in column B
                    total_row = None
                    for row in range(1, copiedSheet.UsedRange.Rows.Count + 1):
                        if copiedSheet.Cells(row, 2).Value == "TOTAL":
                            total_row = row
                            break

                    if total_row is None:
                        raise ValueError("Row with 'TOTAL' not found in column B.")
                    
                    # Define the columns to extract based on sheet name
                    column_indices = [3, 7]

                    # Get values from specified cells in the "TOTAL" row based on sheet name
                    numeric_values = [copiedSheet.Cells(total_row, col).Value for col in column_indices]
                    cell_values = []

                    for value in numeric_values:
                        if isinstance(value, (int, float)):
                            cell_values.append(round(value))
                        else:
                            # Handle non-numeric values, such as strings
                            cell_values.append(value)

                    # Calculate the desired value
                    result_value = ((cell_values[1] - cell_values[0]) / cell_values[0]) * 100

                    # Create a new slide in the PowerPoint presentation
                    slide = presentation.slides[slide_index]

                    # Add data to the slide with individual positions for each textbox:left,top,width,height
                    textbox1 = slide.shapes.add_textbox(Inches(2.5), Inches(1.9), Inches(1), Inches(0.5))
                    textbox1.text_frame.text = f"{round(result_value, 2)}%"  # Round to 2 decimal places 

                else:
                    raise ValueError("Column 'Total' not found in sheet 'AVG. ATTENDANCE'")

                return

        # No sheet found, raise an exception or handle accordingly
        raise ValueError(f"No sheet found: {sheet_name}")

    except Exception as e:
        logging.error(f"Error while processing {excel_path}: {e}")

    finally:
        excelWorkbook.Close(SaveChanges=False)

def process_excel_files_in_folder(folder_path, pptx_template_path):
    xlApp = win32com.client.Dispatch('Excel.Application')
    pptApp = win32com.client.Dispatch('PowerPoint.Application')

    try:
        # Load the PowerPoint template once outside the loop
        template_presentation = Presentation(pptx_template_path)

        for file_name in os.listdir(folder_path):
            if file_name.endswith('.xlsx'):
                # Reuse the same template for each Excel file
                presentation = Presentation(pptx_template_path)

                excel_path = os.path.join(folder_path, file_name)
                excel_file_name = os.path.splitext(file_name)[0]
                """
                # Instructions to copy chart slide
                try:
                    copy_chart_to_existing_slide(excel_path, 'MRR', presentation, 9)
                    copy_chart_to_existing_slide(excel_path, 'CSR DISTRIBUTION', presentation, 13)
                    copy_chart_to_existing_slide(excel_path, 'CHURCH ANALYSIS1', presentation, 23)
                except win32com.client.pywintypes.com_error as e:
                    print(f"Error in recommendation slides for {excel_path}: {e}")
                    logging.error(f"Error in recommendation slides for {excel_path}: {e}")
                """
                # Instructions to copy table to slide
                try:
                    copy_avg_attendance_to_slide(excel_path, 'AVG. ATTENDANCE', presentation, 4, copy_count=16)
                    copy_avg_attendance_to_slide(excel_path, 'AVG. ATTENDANCE', presentation, 5, copy_count=16)
                    copy_table_to_slide(excel_path, 'MRR', presentation, 8, font_size=Pt(15))
                    copy_table_to_slide(excel_path, 'CSR', presentation, 11, font_size=Pt(15))
                    copy_table_to_slide(excel_path, 'CSR DISTRIBUTION', presentation, 12, font_size=Pt(15))
                    copy_table_to_slide(excel_path, 'CONVERTS', presentation, 16, copy_count=16, font_size=Pt(9))
                    copy_table_to_slide(excel_path, 'HOUSE FELLOWSHIP', presentation, 19, copy_count=16, font_size=Pt(9))
                    copy_table_to_slide(excel_path, 'CHURCH ANALYSIS1', presentation, 22, copy_count=16, font_size=Pt(15))
                except win32com.client.pywintypes.com_error as e:
                    print(f"Error in recommendation slides for {excel_path}: {e}")
                    logging.error(f"Error in recommendation slides for {excel_path}: {e}")
              
                # Instructions to copy cells to slide
                try:
                    copy_cells_to_slide(excel_path, 'AVG. ATTENDANCE', presentation, 6)
                    copy_cells_to_slide(excel_path, 'HOUSE FELLOWSHIP', presentation, 20)
                    copy_cells_to_slide(excel_path, 'CONVERTS', presentation, 17)
                    copy_mrrcells_to_slide(excel_path, 'MRR', presentation, 9)
                    copy_csrcells_to_slide(excel_path, 'CSR', presentation, 14)
                    copy_cacells_to_slide(excel_path, 'CHURCH ANALYSIS1', presentation, 24)
                except win32com.client.pywintypes.com_error as e:
                    print(f"Error in recommendation slides for {excel_path}: {e}")
                    logging.error(f"Error in recommendation slides for {excel_path}: {e}")
                  
                # Recommendation pages
                try:
                    recommendation_CA_slide(excel_path, 'CHURCH ANALYSIS1', presentation, 28)
                    recommendation_MRR_slide(excel_path, 'MRR', presentation, 26)
                    recommendation_CSR_slide(excel_path, 'CSR', presentation, 27)
                    recommendation_HF_slide(excel_path, 'HOUSE FELLOWSHIP', presentation, 28)
                    recommendation_CONVERTS_slide(excel_path, 'CONVERTS', presentation, 27)          
                    recommendation_ATTENDANCE_slide(excel_path, 'AVG. ATTENDANCE', presentation, 26)
                except win32com.client.pywintypes.com_error as e:
                    print(f"Error in recommendation slides for {excel_path}: {e}")
                    logging.error(f"Error in recommendation slides for {excel_path}: {e}")
                  
                output_pptx_path = os.path.join(folder_path, f'{excel_file_name}_updated.pptx')
                presentation.save(output_pptx_path)

                logging.info(f'Finished populating {file_name}')
                # Add a sleep interval (adjust as needed)
                time.sleep(15)
    except Exception as e:
        logging.error(f"Error during processing: {e}")

    finally:
        # Quit Excel and PowerPoint after processing all files
        xlApp.DisplayAlerts = False  # Disable Excel alerts to prevent prompts when closing
        xlApp.Quit()
        del xlApp 
        pptApp.Quit()

def main():
    if len(sys.argv) == 3:
        folder_path = sys.argv[1]
        pptx_path = sys.argv[2]
    else:
        # Use default paths if command-line arguments are not provided
        folder_path = r"C:\Users\PMD - FEMI\Desktop\Tiers\Tier 2\cleaned_files"
        pptx_path = r"C:\Users\PMD - FEMI\Desktop\PPTX_TEMPLATES\Q1template.pptx"

    process_excel_files_in_folder(folder_path, pptx_path)

if __name__ == "__main__":
    main()